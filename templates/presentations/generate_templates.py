from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
MASTER_NAME = "Codyssey_B1-B7_Mission_Evaluation_Master_Template.pptx"
CORE_NAME = "Codyssey_Mission_Evaluation_Core_22slides.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x08, 0x1B, 0x2F)
NAVY_2 = RGBColor(0x11, 0x2B, 0x47)
BLUE = RGBColor(0x2F, 0x80, 0xED)
TEAL = RGBColor(0x19, 0xA8, 0x9D)
AMBER = RGBColor(0xF3, 0xA6, 0x1D)
RED = RGBColor(0xD9, 0x5D, 0x58)
GREEN = RGBColor(0x2F, 0xA6, 0x6A)
LIGHT = RGBColor(0xF6, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x22, 0x2D)
MUTED = RGBColor(0x65, 0x72, 0x81)
LINE = RGBColor(0xD9, 0xE1, 0xEA)
PALE_BLUE = RGBColor(0xE9, 0xF2, 0xFF)
PALE_TEAL = RGBColor(0xE9, 0xF8, 0xF6)
PALE_AMBER = RGBColor(0xFF, 0xF5, 0xDF)
PALE_RED = RGBColor(0xFC, 0xED, 0xEC)
PALE_GREEN = RGBColor(0xE9, 0xF7, 0xEF)

FONT = "Noto Sans CJK KR"
FONT_MONO = "Aptos Mono"


def rgb(hexstr: str) -> RGBColor:
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, *, size=18, color=INK, bold=False,
             font=FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.03, line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullets(slide, items: Iterable[str], x, y, w, h, *, size=14, color=INK, gap=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
        p.level = 0
        p.text = "• " + item
    return box


def add_rect(slide, x, y, w, h, fill, *, line=LINE, radius=True, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_card(slide, title, body, x, y, w, h, *, accent=BLUE, pale=None, number=None):
    pale = pale or PALE_BLUE
    add_rect(slide, x, y, w, h, WHITE, line=LINE)
    add_rect(slide, x, y, Inches(0.09), h, accent, line=accent, radius=False)
    if number is not None:
        add_text(slide, str(number), x+Inches(0.24), y+Inches(0.18), Inches(0.48), Inches(0.38), size=11, bold=True, color=accent)
        tx = x+Inches(0.72)
        tw = w-Inches(0.92)
    else:
        tx = x+Inches(0.28)
        tw = w-Inches(0.48)
    add_text(slide, title, tx, y+Inches(0.16), tw, Inches(0.42), size=13.5, bold=True, color=INK)
    add_text(slide, body, x+Inches(0.28), y+Inches(0.68), w-Inches(0.5), h-Inches(0.84), size=10.5, color=MUTED)


def add_header(slide, title, section, *, accent=BLUE, slide_no=None, claim=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), accent, line=accent, radius=False)
    add_text(slide, section.upper(), Inches(0.62), Inches(0.34), Inches(2.5), Inches(0.3), size=9.5, bold=True, color=accent)
    add_text(slide, title, Inches(0.62), Inches(0.68), Inches(11.7), Inches(0.72), size=24, bold=True, color=INK)
    if claim:
        add_text(slide, claim, Inches(0.64), Inches(1.38), Inches(11.6), Inches(0.38), size=11.5, color=MUTED)
    add_text(slide, "CODYSSEY BASIC · MISSION EVALUATION", Inches(0.62), Inches(7.17), Inches(4.8), Inches(0.18), size=7.5, color=MUTED)
    if slide_no is not None:
        add_text(slide, f"{slide_no:02d}", Inches(12.12), Inches(7.10), Inches(0.55), Inches(0.24), size=8.5, bold=True, color=accent, align=PP_ALIGN.RIGHT)


def add_pill(slide, text, x, y, w, *, fill=PALE_BLUE, color=BLUE):
    add_rect(slide, x, y, w, Inches(0.34), fill, line=fill)
    add_text(slide, text, x+Inches(0.06), y+Inches(0.055), w-Inches(0.12), Inches(0.22), size=8.5, bold=True, color=color, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, *, color=BLUE, width=2.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def slide_cover(prs, no, title="[B?-?] 미션 평가 발표", subtitle="Requirement → Design → Verification → Evidence"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    add_rect(s, Inches(0.72), Inches(0.70), Inches(0.12), Inches(5.85), TEAL, line=TEAL, radius=False)
    add_pill(s, "CODYSSEY · AI/SW BASIC", Inches(1.18), Inches(0.84), Inches(2.25), fill=rgb("#163758"), color=rgb("#8BC9FF"))
    add_text(s, title, Inches(1.18), Inches(1.62), Inches(10.8), Inches(1.15), size=31, bold=True, color=WHITE)
    add_text(s, "평가자가 짧은 시간 안에 요구사항 충족과 실제 검증 근거를 판단할 수 있는 발표 구조", Inches(1.2), Inches(2.83), Inches(9.8), Inches(0.7), size=15, color=rgb("#C7D6E6"))
    add_text(s, subtitle, Inches(1.2), Inches(4.25), Inches(8.8), Inches(0.45), size=14, bold=True, color=TEAL)
    add_text(s, "발표자  [이름]   ·   Repository  [URL]   ·   Date  [YYYY-MM-DD]", Inches(1.2), Inches(5.35), Inches(10.3), Inches(0.4), size=10.5, color=rgb("#9DB2C8"))
    add_text(s, f"{no:02d}", Inches(11.85), Inches(6.82), Inches(0.7), Inches(0.25), size=9, bold=True, color=rgb("#7894AE"), align=PP_ALIGN.RIGHT)


def slide_exec(prs, no):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_header(s, "핵심 결론을 먼저 제시한다", "Executive Summary", accent=TEAL, slide_no=no,
               claim="평가자는 첫 60초 안에 ‘무엇을 만들었고, 어떻게 검증했는가’를 이해할 수 있어야 한다.")
    cards = [("MISSION", "[공식 요구 핵심 1문장]", BLUE, PALE_BLUE),("SOLUTION", "[내 구현의 핵심 접근 1문장]", TEAL, PALE_TEAL),("VERIFICATION", "[자동/실환경 검증 결과 요약]", GREEN, PALE_GREEN),("LEARNING", "[내가 설명할 수 있게 된 핵심 개념]", AMBER, PALE_AMBER)]
    x0=Inches(0.72); y=Inches(2.20); gap=Inches(0.22); w=Inches(2.88); h=Inches(2.42)
    for i,(t,b,a,p) in enumerate(cards): add_card(s,t,b,x0+i*(w+gap),y,w,h,accent=a,pale=p,number=i+1)
    add_rect(s, Inches(0.72), Inches(5.08), Inches(11.90), Inches(1.24), WHITE, line=LINE)
    add_text(s, "한 문장 핵심 주장", Inches(0.98), Inches(5.30), Inches(1.55), Inches(0.3), size=10, bold=True, color=TEAL)
    add_text(s, "“[미션의 핵심 요구를 실제 실행과 추적 가능한 증빙으로 충족했다.]”", Inches(2.48), Inches(5.23), Inches(9.65), Inches(0.54), size=16, bold=True, color=INK)


def slide_agenda(prs, no):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_header(s, "발표 논리는 요구사항에서 증빙까지 한 방향으로 흐른다", "Agenda", accent=BLUE, slide_no=no)
    stages=[("01","MISSION","공식 요구"),("02","DESIGN","설계 선택"),("03","BUILD","구현"),("04","VERIFY","검증"),("05","EVIDENCE","증빙"),("06","EXPLAIN","학습·Q&A")]
    x=Inches(0.75); y=Inches(2.25); w=Inches(1.72); gap=Inches(0.25)
    for i,(n,t,b) in enumerate(stages):
        add_rect(s,x+i*(w+gap),y,w,Inches(2.0),WHITE,line=LINE)
        add_text(s,n,x+i*(w+gap)+Inches(0.18),y+Inches(0.20),Inches(0.45),Inches(0.3),size=10,bold=True,color=BLUE)
        add_text(s,t,x+i*(w+gap)+Inches(0.18),y+Inches(0.67),w-Inches(0.36),Inches(0.35),size=13,bold=True,color=INK)
        add_text(s,b,x+i*(w+gap)+Inches(0.18),y+Inches(1.20),w-Inches(0.36),Inches(0.45),size=10,color=MUTED)
        if i<5: add_arrow(s,x+i*(w+gap)+w,y+Inches(1.0),x+(i+1)*(w+gap)-Inches(0.05),y+Inches(1.0),color=TEAL,width=1.6)


def slide_mission(prs, no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_header(s,"공식 미션을 ‘평가 가능한 요구사항’으로 다시 정리한다","Mission Brief",accent=BLUE,slide_no=no,claim="Source of Truth: Mission PDF → Evaluation → 실제 Repository / Evidence")
    add_card(s,"MISSION OBJECTIVE","[학습자가 최종적으로 할 수 있어야 하는 것]",Inches(0.72),Inches(2.02),Inches(3.65),Inches(1.62),accent=BLUE)
    add_card(s,"FINAL OUTPUT","[필수 결과물과 제출물]",Inches(4.57),Inches(2.02),Inches(3.65),Inches(1.62),accent=TEAL,pale=PALE_TEAL)
    add_card(s,"CONSTRAINTS","[제약조건·금지사항·환경 조건]",Inches(8.42),Inches(2.02),Inches(3.65),Inches(1.62),accent=AMBER,pale=PALE_AMBER)
    add_rect(s,Inches(0.72),Inches(4.04),Inches(11.35),Inches(2.0),WHITE,line=LINE)
    add_text(s,"발표에서 반드시 답할 3가지",Inches(0.98),Inches(4.30),Inches(2.1),Inches(0.35),size=12,bold=True,color=INK)
    add_bullets(s,["무엇을 구현해야 했는가?","내 구현은 어디에서 그 요구를 충족하는가?","실제로 동작했다는 근거는 무엇인가?"],Inches(3.02),Inches(4.24),Inches(8.45),Inches(1.25),size=12)


def slide_matrix(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_header(s,"요구사항마다 구현·검증·증빙을 1:1로 연결한다","Requirement → Evidence Matrix",accent=GREEN,slide_no=no)
    x=Inches(0.72); y=Inches(2.0); widths=[1.05,3.0,2.6,2.4,2.3]; headers=["ID","Requirement","Implementation","Verification","Evidence"]
    cx=x
    for w,h in zip(widths,headers):
        add_rect(s,cx,y,Inches(w),Inches(0.58),NAVY,line=NAVY,radius=False); add_text(s,h,cx+Inches(0.07),y+Inches(0.14),Inches(w-0.14),Inches(0.25),size=9.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER); cx+=Inches(w)
    rows=[("R-01","[필수 기능]","[파일/함수]","[테스트 명령]","[로그/캡처]"),("R-02","[제약 조건]","[설정/정책]","[정적/실행 검증]","[설정 증빙]"),("R-03","[오류 처리]","[예외/복구]","[실패 시나리오]","[Before/After]"),("R-04","[문서/설명]","[README/Guide]","[링크 점검]","[문서 위치]"),("R-05","[배포/협업]","[URL/PR]","[외부 접속/리뷰]","[URL/PR 링크]")]
    ry=y+Inches(0.58)
    for ri,row in enumerate(rows):
        cx=x; fill=WHITE if ri%2==0 else rgb("#F2F5F8")
        for w,val in zip(widths,row):
            add_rect(s,cx,ry,Inches(w),Inches(0.66),fill,line=LINE,radius=False); add_text(s,val,cx+Inches(0.07),ry+Inches(0.16),Inches(w-0.14),Inches(0.27),size=8.8,color=INK,align=PP_ALIGN.CENTER); cx+=Inches(w)
        ry+=Inches(0.66)
    add_text(s,"TIP · PASS를 먼저 쓰지 말고 Evidence가 존재하는지부터 확인",Inches(0.82),Inches(6.12),Inches(7.2),Inches(0.3),size=10,bold=True,color=GREEN)


def slide_problem(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_header(s,"문제 → 원인 → 목표를 한 장에서 연결한다","Problem → Objective",accent=AMBER,slide_no=no)
    add_card(s,"PROBLEM","[관찰한 문제 / 공식 미션의 문제 상황]",Inches(0.75),Inches(2.15),Inches(3.45),Inches(2.25),accent=RED,pale=PALE_RED)
    add_arrow(s,Inches(4.25),Inches(3.28),Inches(4.85),Inches(3.28),color=AMBER)
    add_card(s,"ROOT CAUSE / GAP","[왜 문제가 발생하는가 / 무엇이 부족한가]",Inches(4.92),Inches(2.15),Inches(3.45),Inches(2.25),accent=AMBER,pale=PALE_AMBER)
    add_arrow(s,Inches(8.42),Inches(3.28),Inches(9.02),Inches(3.28),color=TEAL)
    add_card(s,"OBJECTIVE","[이번 구현에서 달성할 검증 가능한 목표]",Inches(9.08),Inches(2.15),Inches(3.45),Inches(2.25),accent=TEAL,pale=PALE_TEAL)
    add_text(s,"목표는 ‘좋아 보이는 기능’이 아니라 평가 가능한 상태로 작성한다.",Inches(0.82),Inches(5.08),Inches(9.7),Inches(0.36),size=12,bold=True,color=INK)
    add_text(s,"예: ‘안정성을 높인다’ → ‘실패 입력에서 비정상 종료 없이 명시적 오류를 반환한다’",Inches(0.82),Inches(5.55),Inches(10.6),Inches(0.35),size=10.5,color=MUTED)


def slide_solution(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_header(s,"해결책은 기능 목록이 아니라 전체 사용자·시스템 흐름으로 보여준다","Solution Overview",accent=TEAL,slide_no=no)
    steps=[("INPUT","사용자/이벤트"),("PROCESS","핵심 로직"),("DATA","저장/상태"),("VERIFY","검증/로그"),("OUTPUT","결과/화면")]
    x=Inches(0.88); y=Inches(2.62); w=Inches(2.05); gap=Inches(0.36); colors=[BLUE,TEAL,AMBER,GREEN,NAVY_2]
    for i,((t,b),c) in enumerate(zip(steps,colors)):
        add_rect(s,x+i*(w+gap),y,w,Inches(1.5),WHITE,line=LINE); add_pill(s,t,x+i*(w+gap)+Inches(0.20),y+Inches(0.20),Inches(0.9),fill=rgb("#EDF2F7"),color=c); add_text(s,b,x+i*(w+gap)+Inches(0.20),y+Inches(0.78),w-Inches(0.40),Inches(0.35),size=11,bold=True,color=INK)
        if i<4:add_arrow(s,x+i*(w+gap)+w,y+Inches(0.76),x+(i+1)*(w+gap)-Inches(0.04),y+Inches(0.76),color=TEAL,width=1.8)
    add_rect(s,Inches(0.88),Inches(4.72),Inches(11.4),Inches(1.28),PALE_TEAL,line=PALE_TEAL); add_text(s,"핵심 설계 선택(Why)",Inches(1.12),Inches(4.98),Inches(2.0),Inches(0.3),size=11,bold=True,color=TEAL); add_text(s,"[선택한 구조]를 사용한 이유는 [평가 요구 / 유지보수 / 검증 가능성] 때문이다.",Inches(3.10),Inches(4.91),Inches(8.52),Inches(0.55),size=14,bold=True,color=INK)


def slide_architecture(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_header(s,"컴포넌트의 책임과 데이터 흐름을 동시에 보여준다","Architecture",accent=BLUE,slide_no=no)
    nodes=[("CLIENT","UI / CLI"),("ENTRY","Router / Command"),("LOGIC","Service / Core"),("DATA","Repository / DB"),("EXTERNAL","API / OS / Cloud")]
    x=Inches(0.78); y=Inches(2.35); w=Inches(2.14); gap=Inches(0.28)
    for i,(t,b) in enumerate(nodes):
        c=[NAVY_2,BLUE,TEAL,AMBER,GREEN][i]; add_rect(s,x+i*(w+gap),y,w,Inches(1.72),WHITE,line=LINE); add_text(s,t,x+i*(w+gap)+Inches(0.18),y+Inches(0.22),w-Inches(0.36),Inches(0.28),size=9,bold=True,color=c); add_text(s,b,x+i*(w+gap)+Inches(0.18),y+Inches(0.73),w-Inches(0.36),Inches(0.50),size=12.5,bold=True,color=INK)
        if i<4:add_arrow(s,x+i*(w+gap)+w,y+Inches(0.85),x+(i+1)*(w+gap)-Inches(0.04),y+Inches(0.85),color=BLUE,width=1.7)
    add_bullets(s,["각 박스에는 ‘무엇을 하는가’보다 ‘어떤 책임을 갖는가’를 작성","화살표에는 Request / Response / Event / Data 등 흐름의 의미를 표기","평가 질문: ‘왜 이렇게 분리했나요?’에 바로 답할 수 있도록 구성"],Inches(0.92),Inches(4.60),Inches(11.1),Inches(1.35),size=11.5)


def slide_stack(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"기술 스택은 이름보다 선택 이유와 역할을 중심으로 정리한다","Technology Stack",accent=TEAL,slide_no=no)
    data=[("LANGUAGE","Python / JS / SQL","[왜 이 언어인가]",BLUE),("FRAMEWORK","FastAPI / React","[어떤 책임인가]",TEAL),("DATA","SQLite / PostgreSQL","[어떤 데이터 요구인가]",AMBER),("TOOLS","Git / GitHub / Test","[어떻게 검증·협업하는가]",GREEN)]
    x=Inches(0.80); y=Inches(2.28); w=Inches(2.82); gap=Inches(0.24)
    for i,(a,b,c,col) in enumerate(data): add_card(s,a,b+"\n"+c,x+i*(w+gap),y,w,Inches(2.62),accent=col,number=i+1)
    add_text(s,"Avoid · 로고만 나열하는 ‘기술 스택 슬라이드’",Inches(0.82),Inches(5.45),Inches(4.3),Inches(0.3),size=10,bold=True,color=RED); add_text(s,"Prefer · ‘요구사항 → 기술 선택 → 책임 → 검증 방식’",Inches(5.15),Inches(5.45),Inches(5.6),Inches(0.3),size=10,bold=True,color=GREEN)


def slide_deep(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"핵심 코드 10~18줄만 보여주고 설계 의도를 설명한다","Technical Deep Dive",accent=TEAL,slide_no=no)
    add_rect(s,Inches(0.72),Inches(2.02),Inches(7.40),Inches(4.30),NAVY,line=NAVY)
    code="01  def [core_function](input):\n02      # 핵심 전제 / invariant\n03      validated = validate(input)\n04      result = service.execute(validated)\n05      evidence.log(result)\n06      return result\n\n# [실제 코드 10~18줄 삽입]"
    add_text(s,code,Inches(1.02),Inches(2.38),Inches(6.75),Inches(3.35),size=12,color=rgb("#D9E7F5"),font=FONT_MONO)
    add_card(s,"WHY","[이 로직이 필요한 이유]",Inches(8.38),Inches(2.02),Inches(3.92),Inches(1.22),accent=BLUE); add_card(s,"CONTRACT","[입력·출력·예외 계약]",Inches(8.38),Inches(3.44),Inches(3.92),Inches(1.22),accent=AMBER,pale=PALE_AMBER); add_card(s,"VERIFY","[어떤 테스트로 검증했는가]",Inches(8.38),Inches(4.86),Inches(3.92),Inches(1.46),accent=GREEN,pale=PALE_GREEN)


def slide_flow(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"하나의 핵심 기능을 입력부터 결과까지 추적한다","Core Logic / Data Flow",accent=TEAL,slide_no=no)
    ys=[2.20,3.20,4.20,5.20]; labels=[("1","INPUT","[사용자 입력 / 요청]"),("2","TRANSFORM","[검증·파싱·상태 변경]"),("3","CORE","[핵심 알고리즘 / 비즈니스 로직]"),("4","OUTPUT","[저장·응답·로그]")]
    for i,(n,t,b) in enumerate(labels):
        add_rect(s,Inches(2.0),Inches(ys[i]),Inches(9.15),Inches(0.76),WHITE,line=LINE); add_rect(s,Inches(2.0),Inches(ys[i]),Inches(0.72),Inches(0.76),[BLUE,TEAL,AMBER,GREEN][i],line=[BLUE,TEAL,AMBER,GREEN][i],radius=False); add_text(s,n,Inches(2.18),Inches(ys[i]+0.20),Inches(0.32),Inches(0.25),size=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER); add_text(s,t,Inches(2.95),Inches(ys[i]+0.17),Inches(1.35),Inches(0.28),size=10,bold=True,color=INK); add_text(s,b,Inches(4.30),Inches(ys[i]+0.17),Inches(6.2),Inches(0.3),size=10.8,color=MUTED)
        if i<3:add_arrow(s,Inches(6.55),Inches(ys[i]+0.76),Inches(6.55),Inches(ys[i+1]-0.04),color=TEAL,width=1.6)


def slide_demo(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"데모는 ‘성공 장면’이 아니라 평가 시나리오 순서로 설계한다","Demo Flow",accent=BLUE,slide_no=no)
    flow=[("01","START","실행/접속"),("02","INPUT","핵심 입력"),("03","ACTION","핵심 기능"),("04","FAIL","예외 1개"),("05","RESULT","정상 결과"),("06","EVIDENCE","로그/DB/PR 확인")]
    x=Inches(0.72); y=Inches(2.42); w=Inches(1.78); gap=Inches(0.25)
    for i,(n,t,b) in enumerate(flow):
        add_card(s,t,b,x+i*(w+gap),y,w,Inches(2.15),accent=BLUE if i<3 else (RED if i==3 else GREEN),number=n)
        if i<5:add_arrow(s,x+i*(w+gap)+w,y+Inches(1.05),x+(i+1)*(w+gap)-Inches(0.05),y+Inches(1.05),color=MUTED,width=1.2)
    add_text(s,"데모 실패 대비: 동일 시나리오의 캡처·로그·테스트 결과를 부록에 준비",Inches(0.82),Inches(5.35),Inches(10.3),Inches(0.35),size=11,bold=True,color=AMBER)


def slide_evidence(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"스크린샷마다 ‘무엇을 증명하는가’를 명시한다","Demo Evidence",accent=GREEN,slide_no=no)
    add_rect(s,Inches(0.75),Inches(2.05),Inches(7.25),Inches(4.20),rgb("#E8EDF3"),line=LINE); add_text(s,"[실제 화면 / 터미널 / 로그 / PR 캡처]",Inches(1.25),Inches(3.60),Inches(6.2),Inches(0.45),size=16,bold=True,color=MUTED,align=PP_ALIGN.CENTER); add_card(s,"EVIDENCE CAPTION","이 캡처는 [요구사항 R-??]가 [실제 환경]에서 [정상 결과]임을 증명한다.",Inches(8.30),Inches(2.05),Inches(4.0),Inches(1.80),accent=GREEN,pale=PALE_GREEN); add_card(s,"TRACEABILITY","Repository / Branch / Commit / Runtime / Executed At",Inches(8.30),Inches(4.05),Inches(4.0),Inches(1.40),accent=BLUE); add_text(s,"Secret·Token·개인정보는 캡처 전에 제거",Inches(8.48),Inches(5.72),Inches(3.6),Inches(0.30),size=10,bold=True,color=RED)


def slide_test(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"테스트는 ‘무엇을 실패하게 해봤는가’까지 보여준다","Test Strategy",accent=GREEN,slide_no=no)
    levels=[("STATIC","Syntax / Lint","[정적 오류]",BLUE),("UNIT","Function / Class","[핵심 로직]",TEAL),("INTEGRATION","CLI / HTTP / DB","[연결 흐름]",AMBER),("RUNTIME","실환경 / 외부 접속","[자동화 불가 항목]",GREEN)]
    x=Inches(0.82); y=Inches(2.14); w=Inches(2.78); gap=Inches(0.26)
    for i,(a,b,c,col) in enumerate(levels): add_card(s,a,b+"\n"+c,x+i*(w+gap),y,w,Inches(2.25),accent=col,number=i+1)
    add_rect(s,Inches(0.82),Inches(4.78),Inches(11.45),Inches(1.20),PALE_GREEN,line=PALE_GREEN); add_text(s,"최소 실패 시나리오",Inches(1.10),Inches(5.05),Inches(1.65),Inches(0.3),size=10,bold=True,color=GREEN); add_text(s,"[잘못된 입력] · [없는 데이터] · [권한 없음] · [외부 API 실패] 중 미션에 맞는 항목을 실제로 검증",Inches(2.75),Inches(4.98),Inches(8.95),Inches(0.45),size=12,bold=True,color=INK)


def slide_dashboard(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"검증 상태를 한눈에 읽을 수 있는 대시보드로 요약한다","Verification Dashboard",accent=GREEN,slide_no=no)
    metrics=[("REQUIREMENTS","[12 / 12]","Mapped",BLUE),("AUTOMATED TEST","[28 / 28]","PASS",GREEN),("RUNTIME","[MAC-V / WIN-V]","Actual",TEAL),("BLOCKER / MAJOR","[0 / 0]","Open",AMBER)]
    x=Inches(0.78); y=Inches(2.18); w=Inches(2.82); gap=Inches(0.25)
    for i,(t,v,b,c) in enumerate(metrics): add_rect(s,x+i*(w+gap),y,w,Inches(2.10),WHITE,line=LINE); add_text(s,t,x+i*(w+gap)+Inches(0.22),y+Inches(0.22),w-Inches(0.44),Inches(0.28),size=9,bold=True,color=MUTED); add_text(s,v,x+i*(w+gap)+Inches(0.22),y+Inches(0.72),w-Inches(0.44),Inches(0.55),size=24,bold=True,color=c); add_text(s,b,x+i*(w+gap)+Inches(0.22),y+Inches(1.48),w-Inches(0.44),Inches(0.3),size=10,bold=True,color=INK)
    add_bullets(s,["IMPLEMENTED ≠ TESTED ≠ PASS — 상태를 섞지 않는다","NEEDS-RUNTIME은 실제 환경 검증 전까지 별도 표시","예상 출력과 실제 출력은 명확히 구분"],Inches(0.98),Inches(4.75),Inches(10.8),Inches(1.25),size=11.5)


def slide_before_after(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"변경의 효과는 Before & After 또는 정량 결과로 증명한다","Results / Before & After",accent=GREEN,slide_no=no)
    add_rect(s,Inches(0.78),Inches(2.05),Inches(5.55),Inches(3.65),PALE_RED,line=rgb("#F0C9C7")); add_text(s,"BEFORE",Inches(1.05),Inches(2.30),Inches(1.2),Inches(0.3),size=11,bold=True,color=RED); add_bullets(s,["[문제 상태 / 수치]","[실패 로그 / 사용자 증상]","[원인 가설]"],Inches(1.05),Inches(3.00),Inches(4.65),Inches(1.45),size=13)
    add_rect(s,Inches(6.70),Inches(2.05),Inches(5.55),Inches(3.65),PALE_GREEN,line=rgb("#C4E7D1")); add_text(s,"AFTER",Inches(6.98),Inches(2.30),Inches(1.2),Inches(0.3),size=11,bold=True,color=GREEN); add_bullets(s,["[조치 후 상태 / 수치]","[성공 로그 / 정상 결과]","[남은 한계]"],Inches(6.98),Inches(3.00),Inches(4.65),Inches(1.45),size=13); add_arrow(s,Inches(6.15),Inches(3.85),Inches(6.70),Inches(3.85),color=TEAL,width=2.4)


def slide_trouble(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"문제 해결은 증상 → 가설 → 검증 → 조치 → 재발방지로 기록한다","Troubleshooting",accent=AMBER,slide_no=no)
    steps=[("SYMPTOM","[관측된 현상]",RED),("HYPOTHESIS","[원인 가설]",AMBER),("VERIFY","[로그/명령 검증]",BLUE),("ACTION","[최소 조치]",TEAL),("PREVENT","[재발 방지]",GREEN)]
    x=Inches(0.76); y=Inches(2.34); w=Inches(2.10); gap=Inches(0.30)
    for i,(t,b,c) in enumerate(steps): add_card(s,t,b,x+i*(w+gap),y,w,Inches(2.25),accent=c,number=i+1)
    add_text(s,"평가 질문 대비: “처음부터 답을 알고 있었나요?” → 증거를 통해 가설을 좁힌 과정을 설명",Inches(0.86),Inches(5.25),Inches(11.0),Inches(0.38),size=10.5,bold=True,color=INK)


def slide_security(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"보안과 신뢰성은 ‘막았다’가 아니라 경계와 실패 처리를 보여준다","Security / Reliability",accent=RED,slide_no=no)
    add_card(s,"IDENTITY","Authentication\n[누구인가]",Inches(0.80),Inches(2.15),Inches(2.75),Inches(2.15),accent=BLUE); add_card(s,"ACCESS","Authorization / ACL\n[무엇을 할 수 있나]",Inches(3.80),Inches(2.15),Inches(2.75),Inches(2.15),accent=AMBER,pale=PALE_AMBER); add_card(s,"SECRET","API Key / Token\n[노출 방지]",Inches(6.80),Inches(2.15),Inches(2.75),Inches(2.15),accent=RED,pale=PALE_RED); add_card(s,"FAILURE","Timeout / Error\n[비정상 종료 방지]",Inches(9.80),Inches(2.15),Inches(2.75),Inches(2.15),accent=GREEN,pale=PALE_GREEN); add_text(s,"검증 예: 비로그인 접근 · 잘못된 권한 · Key 미설정 · API timeout",Inches(0.92),Inches(5.12),Inches(10.5),Inches(0.35),size=11,bold=True,color=RED)


def slide_git(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"협업은 결과 파일보다 Issue·PR·Review의 추적 가능한 흐름으로 증명한다","Collaboration / Git",accent=BLUE,slide_no=no)
    nodes=[("ISSUE","작업 정의"),("BRANCH","격리 작업"),("COMMIT","의미 있는 변경"),("PR","What / Why / How"),("REVIEW","피드백 반영"),("MERGE","main 통합")]
    x=Inches(0.72); y=Inches(2.55); w=Inches(1.76); gap=Inches(0.25)
    for i,(t,b) in enumerate(nodes): add_card(s,t,b,x+i*(w+gap),y,w,Inches(1.78),accent=BLUE if i<4 else GREEN,number=i+1)
    add_text(s,"필수 증빙 예: PR 링크 · 리뷰 코멘트 · 충돌 해결 기록 · git log --graph",Inches(0.86),Inches(5.10),Inches(10.4),Inches(0.35),size=11,bold=True,color=INK)


def slide_learning(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"학습성과는 ‘사용했다’가 아니라 코드 근거로 설명할 수 있는가로 평가한다","Learning Outcomes",accent=TEAL,slide_no=no)
    prompts=[("CONCEPT","[핵심 개념]을 내 코드로 설명할 수 있다."),("WHY","왜 이 구조/자료구조/패턴을 선택했는지 말할 수 있다."),("FLOW","입력 → 내부 처리 → 결과 흐름을 추적할 수 있다."),("FAILURE","실패 상황과 복구/예외 처리를 설명할 수 있다.")]
    x=Inches(0.84); y=Inches(2.15); w=Inches(5.55); h=Inches(1.48); gapx=Inches(0.35); gapy=Inches(0.28)
    for i,(t,b) in enumerate(prompts): col=i%2; row=i//2; add_card(s,t,b,x+col*(w+gapx),y+row*(h+gapy),w,h,accent=[BLUE,TEAL,AMBER,GREEN][i],number=i+1)
    add_text(s,"발표자는 ‘정의’보다 ‘내 구현에서 어디에 있고 어떻게 동작하는가’를 우선 설명",Inches(0.92),Inches(5.62),Inches(11.0),Inches(0.35),size=10.5,bold=True,color=TEAL)


def slide_qa(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"평가 질문은 Why·How·Evidence의 세 축으로 준비한다","Evaluation Q&A",accent=AMBER,slide_no=no)
    qs=[("WHY","왜 이 방법을 선택했나요?","[대안 비교 + 요구사항 연결]"),("HOW","실제로 어떻게 동작하나요?","[코드/흐름 근거]"),("EVIDENCE","동작했다는 근거는 무엇인가요?","[테스트/로그/URL/PR]"),("LIMIT","어떤 한계가 있나요?","[현재 범위와 고도화 분리]")]
    x=Inches(0.84); y=Inches(2.08); w=Inches(5.55); h=Inches(1.62); gap=Inches(0.33)
    for i,(t,q,a) in enumerate(qs):
        col=i%2; row=i//2; xx=x+col*(w+gap); yy=y+row*(h+gap); add_rect(s,xx,yy,w,h,WHITE,line=LINE); add_pill(s,t,xx+Inches(0.20),yy+Inches(0.18),Inches(0.78),fill=PALE_AMBER,color=AMBER); add_text(s,q,xx+Inches(1.12),yy+Inches(0.20),w-Inches(1.36),Inches(0.35),size=11.5,bold=True,color=INK); add_text(s,a,xx+Inches(1.12),yy+Inches(0.78),w-Inches(1.36),Inches(0.4),size=10,color=MUTED)


def slide_limit(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"현재 미션 통과 범위와 후속 고도화를 명확히 분리한다","Limitations / Next",accent=AMBER,slide_no=no)
    add_rect(s,Inches(0.82),Inches(2.14),Inches(5.40),Inches(3.45),WHITE,line=LINE); add_text(s,"CURRENT LIMITATIONS",Inches(1.12),Inches(2.42),Inches(2.35),Inches(0.3),size=10,bold=True,color=RED); add_bullets(s,["[공식 범위 밖이라 구현하지 않은 항목]","[현재 데이터/환경의 제약]","[검증에서 남은 실제 리스크]"],Inches(1.12),Inches(3.05),Inches(4.60),Inches(1.35),size=12)
    add_rect(s,Inches(6.55),Inches(2.14),Inches(5.40),Inches(3.45),PALE_TEAL,line=rgb("#C7EAE6")); add_text(s,"NEXT — ADVANCED / PRO",Inches(6.86),Inches(2.42),Inches(2.35),Inches(0.3),size=10,bold=True,color=TEAL); add_bullets(s,["[성능 / 보안 / 확장성 고도화]","[추가 테스트 / CI / 관측성]","[연구·포트폴리오 확장]"],Inches(6.86),Inches(3.05),Inches(4.60),Inches(1.35),size=12)


def slide_conclusion(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s,NAVY); add_pill(s,"CONCLUSION",Inches(0.95),Inches(0.85),Inches(1.45),fill=rgb("#163758"),color=rgb("#8BC9FF")); add_text(s,"요구사항을 구현했고, 실제 검증과 증빙으로 설명할 수 있다",Inches(0.95),Inches(1.55),Inches(11.2),Inches(1.15),size=28,bold=True,color=WHITE); add_text(s,"[미션 핵심 결론 1문장]",Inches(0.98),Inches(3.05),Inches(10.0),Inches(0.55),size=17,bold=True,color=TEAL); add_bullets(s,["Requirement — [핵심 요구]","Implementation — [핵심 구현]","Verification — [핵심 검증]","Evidence — [핵심 증빙]"],Inches(1.00),Inches(4.05),Inches(6.2),Inches(1.45),size=13,color=rgb("#D2DFEC")); add_text(s,"Q&A",Inches(9.55),Inches(4.25),Inches(2.0),Inches(0.65),size=25,bold=True,color=WHITE,align=PP_ALIGN.CENTER); add_text(s,f"{no:02d}",Inches(11.95),Inches(6.85),Inches(0.55),Inches(0.22),size=8,bold=True,color=rgb("#7894AE"),align=PP_ALIGN.RIGHT)


def slide_domain(prs,no,title,accent,left_title,left_items,right_title,right_items,flow_items=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,title,"Domain-specific layout",accent=accent,slide_no=no)
    add_rect(s,Inches(0.78),Inches(2.02),Inches(5.55),Inches(3.80),WHITE,line=LINE); add_text(s,left_title,Inches(1.06),Inches(2.30),Inches(2.5),Inches(0.35),size=11,bold=True,color=accent); add_bullets(s,left_items,Inches(1.05),Inches(2.95),Inches(4.7),Inches(2.10),size=11.5)
    add_rect(s,Inches(6.62),Inches(2.02),Inches(5.55),Inches(3.80),rgb("#F0F4F8"),line=LINE); add_text(s,right_title,Inches(6.90),Inches(2.30),Inches(2.6),Inches(0.35),size=11,bold=True,color=accent); add_bullets(s,right_items,Inches(6.90),Inches(2.95),Inches(4.7),Inches(2.10),size=11.5)
    if flow_items:
        y=Inches(5.22); x=Inches(0.98); total=len(flow_items); w=Inches(10.95/total)
        for i,item in enumerate(flow_items): add_pill(s,item,x+i*w,y,w-Inches(0.08),fill=rgb("#E9EEF5"),color=accent)


def slide_academic(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"연구·학회 발표에서는 문제의 Gap과 검증 설계를 더 깊게 보여준다","Academic / Research Framing",accent=BLUE,slide_no=no)
    frames=[("PROBLEM / GAP","기존 접근의 한계\n[왜 필요한가]",RED),("METHOD","제안 방법 / 구현\n[무엇을 했나]",BLUE),("EVALUATION","실험·검증 설계\n[어떻게 확인했나]",TEAL),("RESULT","정량·정성 결과\n[무엇이 달라졌나]",GREEN),("LIMITATION","일반화·제약\n[어디까지 말할 수 있나]",AMBER)]
    x=Inches(0.72); y=Inches(2.34); w=Inches(2.18); gap=Inches(0.25)
    for i,(t,b,c) in enumerate(frames): add_card(s,t,b,x+i*(w+gap),y,w,Inches(2.45),accent=c,number=i+1)
    add_text(s,"학부/미션 평가의 근거를 유지하되, 대안 비교·재현성·일반화 가능성을 추가",Inches(0.84),Inches(5.40),Inches(10.8),Inches(0.35),size=10.5,bold=True,color=BLUE)


def slide_style(prs,no):
    s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_header(s,"모든 미션 발표에 동일한 시각 언어와 작성 규칙을 사용한다","Style Guide / Master Prompt",accent=TEAL,slide_no=no)
    colors=[("NAVY","#081B2F",NAVY),("BLUE","#2F80ED",BLUE),("TEAL","#19A89D",TEAL),("AMBER","#F3A61D",AMBER),("RED","#D95D58",RED),("GREEN","#2FA66A",GREEN)]
    x=Inches(0.80); y=Inches(2.05)
    for i,(n,h,c) in enumerate(colors):
        xx=x+i*Inches(1.93); add_rect(s,xx,y,Inches(1.66),Inches(0.72),c,line=c); add_text(s,n,xx,y+Inches(0.16),Inches(1.66),Inches(0.22),size=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER); add_text(s,h,xx,y+Inches(0.86),Inches(1.66),Inches(0.22),size=8,color=MUTED,align=PP_ALIGN.CENTER)
    add_card(s,"TITLE RULE","제목만 읽어도 발표 논리가 이어지는 ‘주장형 제목’을 사용",Inches(0.82),Inches(3.45),Inches(5.45),Inches(1.22),accent=BLUE); add_card(s,"EVIDENCE RULE","PASS에는 실제 실행·테스트·로그·URL·PR 등 추적 가능한 근거를 연결",Inches(6.58),Inches(3.45),Inches(5.45),Inches(1.22),accent=GREEN,pale=PALE_GREEN); add_card(s,"LANGUAGE RULE","핵심 영문 용어는 한글(English, 약어)로 첫 등장에 병기",Inches(0.82),Inches(4.90),Inches(5.45),Inches(1.22),accent=TEAL,pale=PALE_TEAL); add_card(s,"PROMPT RULE","Mission → Repository → Test → Evidence를 입력으로 사용하고 일반 지식이 공식 요구를 바꾸지 않게 함",Inches(6.58),Inches(4.90),Inches(5.45),Inches(1.22),accent=AMBER,pale=PALE_AMBER)


def build_deck(path: Path, core=False):
    prs=Presentation(); prs.slide_width=SLIDE_W; prs.slide_height=SLIDE_H
    funcs=[slide_cover,slide_exec,slide_agenda,slide_mission,slide_matrix,slide_problem,slide_solution,slide_architecture,slide_stack,slide_deep,slide_flow,slide_demo,slide_evidence,slide_test,slide_dashboard,slide_before_after,slide_trouble,slide_security,slide_git,slide_learning,slide_qa,slide_limit,slide_conclusion]
    selected={1,2,4,5,6,7,8,10,11,12,13,14,15,16,17,18,19,20,21,22,23,34} if core else set(range(1,35))
    no=0
    for idx,f in enumerate(funcs,start=1):
        if idx not in selected: continue
        no+=1; f(prs,no)
    domain_specs=[
        (24,"Linux/OS · 관제·권한·장애 분석",BLUE,"필수 시각화",["시스템 구성도 / 권한표","monitor.sh / cron 흐름","로그 타임라인 / Before & After"],"핵심 평가 질문",["왜 최소 권한이 필요한가?","장애 원인을 어떤 증거로 좁혔나?","실행 환경과 Evidence는 추적 가능한가?"],["USER","SSH/UFW","APP","MONITOR","LOG"]),
        (25,"Python/CLI · 유지보수 가능한 작은 서비스",TEAL,"필수 시각화",["CLI 명령 흐름","Model–Repository–Service–CLI","Generator / Decorator / Type Hint"],"핵심 평가 질문",["왜 스트리밍이 필요한가?","공통 관심사를 왜 분리했나?","타입 계약이 어떤 오류를 줄이는가?"],["CLI","VALIDATE","SERVICE","STORE","REPORT"]),
        (26,"Git 협업 · Issue에서 Merge까지",BLUE,"필수 시각화",["Issue→Branch→PR→Review→Merge","충돌 마커 / 해결 과정","팀원별 PR·Review 증빙"],"핵심 평가 질문",["왜 GitHub Flow인가?","충돌은 왜 발생했나?","리뷰 반영 흔적은 어디에 있나?"],["ISSUE","BRANCH","PR","REVIEW","MERGE"]),
        (27,"자료구조·알고리즘 · 내부 동작을 시각화",AMBER,"필수 시각화",["노드·포인터·버킷 구조","연산 전/후 상태","시간복잡도와 불변조건"],"핵심 평가 질문",["왜 O(1)인가?","어떤 자료구조가 병목을 해결하나?","최악/평균 복잡도는?"],["INPUT","STRUCTURE","OPERATION","COMPLEXITY"]),
        (28,"Web/Frontend · Event → State → Render",TEAL,"필수 시각화",["반응형 화면 비교","Component tree / DOM","Loading·Error·Empty 상태"],"핵심 평가 질문",["상태는 어디에 두었나?","리렌더링은 언제 일어나나?","실패 UI는 일관적인가?"],["EVENT","STATE","FETCH","RENDER"]),
        (29,"DB/Backend · 관계와 요청 흐름",BLUE,"필수 시각화",["ERD / PK·FK","Request→Router→Service→Repository→DB","인증·인가 / 상태 변경"],"핵심 평가 질문",["왜 이 관계인가?","무결성은 어떻게 보장하나?","비즈니스 로직은 어느 계층인가?"],["REQUEST","ROUTER","SERVICE","REPOSITORY","DB"]),
        (30,"Cloud · 네트워크 경계와 최소 권한",AMBER,"필수 시각화",["VPC/Subnet/Route/IGW/EC2","외부→서비스 트래픽 흐름","Security Group / IAM 경계"],"핵심 평가 질문",["외부 요청이 어떻게 도달하나?","SG와 IAM 차이는?","과금 리소스는 정리했나?"],["INTERNET","IGW","SUBNET","EC2","APP"]),
        (31,"AI API · Prompt → API → Validation",TEAL,"필수 시각화",["git diff 입력","Prompt 구조","API 파라미터 / 오류 처리"],"핵심 평가 질문",["결과 품질을 어떻게 제어했나?","Key는 어떻게 보호했나?","실패 시 어떻게 복구하나?"],["DIFF","PROMPT","API","VALIDATE","OUTPUT"]),
        (32,"Term Project · End-to-End 서비스",GREEN,"필수 시각화",["사용자→Web→API→AI→DB→Cloud","ERD / API 명세","팀 역할 / PR / 배포 URL"],"핵심 평가 질문",["사용자별 데이터는 어떻게 분리되나?","AI 실패 시 서비스는?","배포된 전체 흐름이 동작하나?"],["USER","WEB","API","AI","DB","CLOUD"]),
    ]
    for idx,title,accent,lt,li,rt,ri,flow in domain_specs:
        if idx not in selected: continue
        no+=1; slide_domain(prs,no,title,accent,lt,li,rt,ri,flow)
    if 33 in selected: no+=1; slide_academic(prs,no)
    if 34 in selected: no+=1; slide_style(prs,no)
    prs.save(path)
    return no


if __name__ == "__main__":
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    master_count=build_deck(OUT_DIR/MASTER_NAME, core=False)
    core_count=build_deck(OUT_DIR/CORE_NAME, core=True)
    print(f"generated {MASTER_NAME}: {master_count} slides")
    print(f"generated {CORE_NAME}: {core_count} slides")
